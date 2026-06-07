from pathlib import Path
from uuid import uuid4

from docx import Document as DocxDocument
from fastapi import HTTPException, UploadFile
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy.orm import Session

from app.config import Settings
from app.models import Document, Project


ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".jpg", ".jpeg", ".png", ".xlsx"}
CATEGORY_BY_TYPE = {
    "pdf": "literature",
    "pptx": "existing_ppt",
    "docx": "case_material",
    "image": "image",
    "excel": "data",
}


def secure_filename(filename: str) -> str:
    allowed = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789._-"
    cleaned = "".join(char if char in allowed else "_" for char in filename.strip())
    return cleaned.strip("._")[:180]


def detect_file_type(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".pptx":
        return "pptx"
    if suffix == ".docx":
        return "docx"
    if suffix in {".jpg", ".jpeg", ".png"}:
        return "image"
    if suffix == ".xlsx":
        return "excel"
    return "other"


def parse_text(path: Path, file_type: str) -> str:
    try:
        if file_type == "pdf":
            reader = PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)[:20000]
        if file_type == "docx":
            doc = DocxDocument(str(path))
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)[:20000]
        if file_type == "pptx":
            prs = Presentation(str(path))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)[:20000]
    except Exception:
        return ""
    return ""


async def save_upload(db: Session, settings: Settings, project_id: int, file: UploadFile) -> Document:
    project = db.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    original = file.filename or "upload.bin"
    suffix = Path(original).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="暂不支持该文件类型")

    content = await file.read()
    max_bytes = settings.max_upload_size_mb * 1024 * 1024
    if len(content) > max_bytes:
        raise HTTPException(status_code=413, detail=f"文件不能超过 {settings.max_upload_size_mb}MB")

    safe_original = secure_filename(original) or f"upload{suffix}"
    filename = f"{uuid4().hex}{suffix}"
    project_dir = settings.upload_path / str(project_id)
    project_dir.mkdir(parents=True, exist_ok=True)
    saved_path = project_dir / filename
    saved_path.write_bytes(content)

    file_type = detect_file_type(original)
    document = Document(
        project_id=project_id,
        filename=filename,
        original_filename=safe_original,
        file_type=file_type,
        file_path=str(saved_path),
        file_size=len(content),
        parsed_text=parse_text(saved_path, file_type),
        document_category=CATEGORY_BY_TYPE.get(file_type, "other"),
    )
    db.add(document)
    project.status = "uploaded"
    db.commit()
    db.refresh(document)
    return document
