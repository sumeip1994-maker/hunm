from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import desc, select
from sqlalchemy.orm import Session

from app.config import Settings
from app.models import Artifact, Project, Slide


FONT_NAME = "Microsoft YaHei"


class PPTService:
    def generate(self, db: Session, settings: Settings, project: Project) -> dict[str, Any]:
        artifacts = self._latest_artifacts(db, project.id)
        outline = artifacts.get("outline")
        sections = ["标题页", "PPT背景", "核心问题", "资料与证据整理", "总结"]
        if outline:
            sections = [item.get("title", "未命名章节") for item in outline.content_json.get("sections", [])]

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._title_slide(prs, project)
        self._agenda_slide(prs, sections)
        for title in sections[1:6]:
            self._content_slide(prs, title, self._bullets_for_section(title, project, artifacts))
        self._summary_slide(prs, artifacts)

        output_path = Path(settings.output_path) / f"project_{project.id}_presentation.pptx"
        prs.save(output_path)

        db.query(Slide).filter(Slide.project_id == project.id).delete()
        slide_titles = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")][: len(prs.slides)]
        for idx, title in enumerate(slide_titles, start=1):
            db.add(
                Slide(
                    project_id=project.id,
                    slide_no=idx,
                    slide_type="title" if idx == 1 else "summary" if idx == len(slide_titles) else "agenda" if idx == 2 else "background",
                    title=title.splitlines()[0] if title else f"第 {idx} 页",
                    content_json={"source": "ai_assisted_ppt"},
                    layout="widescreen",
                )
            )
        project.status = "ppt_ready"
        db.commit()
        return {"download_url": f"/projects/{project.id}/ppt/download", "filename": output_path.name}

    def _latest_artifacts(self, db: Session, project_id: int) -> dict[str, Artifact]:
        rows = db.scalars(
            select(Artifact)
            .where(Artifact.project_id == project_id)
            .order_by(Artifact.type, desc(Artifact.version))
        ).all()
        artifacts: dict[str, Artifact] = {}
        for artifact in rows:
            if artifact.type not in artifacts:
                artifacts[artifact.type] = artifact
        return artifacts

    def _add_title(self, slide, text: str, top: float = 0.55) -> None:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.8), Inches(0.7))
        frame = box.text_frame
        frame.text = text
        paragraph = frame.paragraphs[0]
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True

    def _title_slide(self, prs: Presentation, project: Project) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, project.title, 2.1)
        subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.6), Inches(1.2))
        text = subtitle.text_frame
        text.text = f"{project.presentation_type} | {date.today().isoformat()}"
        text.paragraphs[0].font.name = FONT_NAME
        text.paragraphs[0].font.size = Pt(20)
        text.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _agenda_slide(self, prs: Presentation, sections: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "目录")
        body = slide.shapes.add_textbox(Inches(1.1), Inches(1.55), Inches(11), Inches(5.2))
        frame = body.text_frame
        for index, section in enumerate(sections[:8], start=1):
            p = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
            p.text = f"{index}. {section}"
            p.font.name = FONT_NAME
            p.font.size = Pt(20)

    def _content_slide(self, prs: Presentation, title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, title)
        body = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(11.5), Inches(4.8))
        frame = body.text_frame
        for text in bullets[:5]:
            p = frame.paragraphs[0] if not frame.text else frame.add_paragraph()
            p.text = text
            p.font.name = FONT_NAME
            p.font.size = Pt(18)

    def _summary_slide(self, prs: Presentation, artifacts: dict[str, Artifact]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "总结")
        body = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(11.5), Inches(4.8))
        frame = body.text_frame
        analysis = artifacts.get("analysis_report").content_json if artifacts.get("analysis_report") else {}
        takeaways = self._list(analysis.get("suggested_focus"))[:2]
        risk_notes = self._list(analysis.get("risk_notes"))[:1]
        bullets = ["Key Takeaways", *takeaways, *risk_notes, "后续可根据审稿建议补充证据来源和优化页面表达"]
        for text in bullets[:5]:
            p = frame.paragraphs[0] if not frame.text else frame.add_paragraph()
            p.text = text
            p.font.name = FONT_NAME
            p.font.size = Pt(22 if text == "Key Takeaways" else 20)

    def _bullets_for_section(self, title: str, project: Project, artifacts: dict[str, Artifact]) -> list[str]:
        analysis = artifacts.get("analysis_report").content_json if artifacts.get("analysis_report") else {}
        directions = artifacts.get("direction_recommendation").content_json if artifacts.get("direction_recommendation") else {}

        summary = str(analysis.get("summary") or "")
        key_questions = self._list(analysis.get("key_questions"))
        focus = self._list(analysis.get("suggested_focus"))
        risks = self._list(analysis.get("risk_notes"))
        direction_items = directions.get("directions") if isinstance(directions.get("directions"), list) else []
        recommended = next((item for item in direction_items if isinstance(item, dict) and item.get("recommended")), None)
        structure = self._list(recommended.get("structure") if isinstance(recommended, dict) else None)

        if "背景" in title or "问题" in title:
            return self._compact([f"核心问题：{project.core_question}", summary, *key_questions[:2]])
        if "资料" in title or "证据" in title or "发现" in title:
            return self._compact([*focus[:4], *structure[:2]])
        if "讨论" in title or "局限" in title or "争议" in title:
            return self._compact([*risks[:3], "讨论时区分证据支持内容、仍不确定内容和下一步补充方向"])
        if "结论" in title or "总结" in title:
            return self._compact([*focus[:3], *risks[:1]])
        return self._compact([summary, *focus[:3], *key_questions[:1]])

    def _compact(self, values: list[str]) -> list[str]:
        cleaned = [value.strip() for value in values if value and value.strip()]
        if cleaned:
            return cleaned
        return ["本页基于已上传资料和AI分析结果生成。", "请补充证据来源、图表和必要的讲者备注。", "避免超出资料证据范围作出诊疗结论。"]

    def _list(self, value: Any) -> list[str]:
        if isinstance(value, list):
            return [str(item) for item in value if str(item).strip()]
        if value:
            return [str(value)]
        return []
